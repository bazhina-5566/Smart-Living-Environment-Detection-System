"""
Create a PPTX with 4 slides, each containing one of the UI screenshots.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image

# Use absolute Windows paths
img_paths = [
    r"E:\暑期信盈达实习\3.0\2023280520毛晋耀显示器三个页面触摸功能加图片轮换1\2023280520毛晋耀显示器三个页面触摸功能加图片轮换\USER\page1_home.jpg",
    r"E:\暑期信盈达实习\3.0\2023280520毛晋耀显示器三个页面触摸功能加图片轮换1\2023280520毛晋耀显示器三个页面触摸功能加图片轮换\USER\page2_control.jpg",
    r"E:\暑期信盈达实习\3.0\2023280520毛晋耀显示器三个页面触摸功能加图片轮换1\2023280520毛晋耀显示器三个页面触摸功能加图片轮换\USER\page3_status.jpg",
    r"E:\暑期信盈达实习\3.0\2023280520毛晋耀显示器三个页面触摸功能加图片轮换1\2023280520毛晋耀显示器三个页面触摸功能加图片轮换\USER\page4_sensor.jpg",
]
titles = [
    "界面一：首页（Page_Home）— 个人主页显示",
    "界面二：外设控制页（Page_One）— LED/FAN触摸按钮",
    "界面三：系统状态页（Page_Two）— 风扇进度条与RGB状态",
    "界面四：传感器数据页（Page_Sensor）— 火焰+光敏实时显示",
]
descs = [
    "全屏240×320 RGB565图片显示，右下角红色文字\"宾哥真帅\"签名",
    "蓝色LED OFF按钮 + 黄色FAN ON按钮，触摸区域(10-110,40-120)和(130-230,40-120)，200ms防抖",
    "黑底白字System Status标题，LED:OFF(红) Fan:ON(绿)，风扇转速40%进度条，RGB绿灯状态显示",
    "Sensor_data标题，火焰传感器0.2%、光敏电阻30.2%实时百分比显示",
]

# Create presentation with widescreen 16:9 ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Blank layout
blank_layout = prs.slide_layouts[6]

for i, (img_path, title, desc) in enumerate(zip(img_paths, titles, descs)):
    slide = prs.slides.add_slide(blank_layout)

    # Add title at top
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True

    # Add description below title
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(14)
    run2.font.italic = True

    # Get image dimensions
    with Image.open(img_path) as im:
        w_px, h_px = im.size

    # Target display area: ~5 inches tall, max ~4.5 inches wide for portrait phone shape
    target_h = 5.5
    target_w = target_h * (w_px / h_px)

    # Make sure it fits within slide width
    if target_w > 8:
        target_w = 8
        target_h = target_w * (h_px / w_px)

    # Center the image horizontally
    left = Inches((13.333 - target_w) / 2)
    top = Inches(1.7)

    slide.shapes.add_picture(img_path, left, top, width=Inches(target_w), height=Inches(target_h))

output_path = r"E:\暑期信盈达实习\3.0\2023280520毛晋耀显示器三个页面触摸功能加图片轮换1\2023280520毛晋耀显示器三个页面触摸功能加图片轮换\USER\ui_screens.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")