"""
合并两个PPT：将4张UI截图页面附加到主PPT后面。
使用方法：
1. 先在金山文档中分别下载两个PPT到本地：
   - 主PPT: https://www.kdocs.cn/l/cnl5gy1uV7J8 → main.pptx
   - UI截图PPT: https://www.kdocs.cn/l/cekrX46ztnU9 → ui_screens.pptx
2. 把两个文件放到本脚本同目录
3. 运行: python merge_ppts.py
4. 输出 combined.pptx 是合并后的PPT
"""
from pptx import Presentation
import os

script_dir = os.path.dirname(os.path.abspath(__file__))
main_path = os.path.join(script_dir, "main.pptx")
ui_path = os.path.join(script_dir, "ui_screens.pptx")
out_path = os.path.join(script_dir, "combined.pptx")

# 打开主PPT
prs_main = Presentation(main_path)
print(f"主PPT: {len(prs_main.slides)} 页")

# 打开UI截图PPT
prs_ui = Presentation(ui_path)
print(f"UI截图PPT: {len(prs_ui.slides)} 页")

# 复制UI页面到主PPT末尾
for i, slide in enumerate(prs_ui.slides):
    # python-pptx不支持直接跨PPT复制slide, 需要手动复制每个shape
    blank = prs_main.slides.add_slide(prs_main.slide_layouts[6])
    for shape in slide.shapes:
        if shape.has_text_frame:
            txBox = blank.shapes.add_textbox(
                shape.left or 0, shape.top or 0,
                shape.width or 0, shape.height or 0
            )
            tf = txBox.text_frame
            tf.text = shape.text_frame.text
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = run.font.size
                    run.font.bold = run.font.bold
        elif shape.shape_type == 13:  # picture
            # 提取图片字节并重新插入
            try:
                img_stream = shape.image.blob
                pic = blank.shapes.add_picture(
                    shape.image.filename if hasattr(shape.image, 'filename') else None,
                    shape.left, shape.top, shape.width, shape.height,
                )
            except Exception:
                # 用bytes方式重新插入
                from io import BytesIO
                img_stream = BytesIO(shape.image.blob)
                blank.shapes.add_picture(
                    img_stream, shape.left, shape.top,
                    shape.width, shape.height
                )
    print(f"  已合并第 {i+1} 页")

prs_main.save(out_path)
print(f"\n合并完成！输出: {out_path}")
print(f"总页数: {len(prs_main.slides)}")